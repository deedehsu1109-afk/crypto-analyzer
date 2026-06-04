#!/usr/bin/env python3
"""
文件內容提取工具
支援：PDF、DOCX、XLSX、PPTX、ODT、ODS、ODP、圖片（PNG/JPG/BMP/TIFF/WEBP）
用法：python extract.py <檔案路徑> [--meta] [--pages X-Y]
"""
from __future__ import annotations
import sys
import os
import json

# 強制 stdout 使用 UTF-8，避免 CP950/Big5 無法編碼特殊字元
if hasattr(sys.stdout, "reconfigure"):
    try:
        sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    except Exception:
        pass
import argparse
import datetime

def _check_dep(pkg: str) -> bool:
    import importlib
    try:
        importlib.import_module(pkg)
        return True
    except ImportError:
        return False


# ── PDF ──────────────────────────────────────────────────────────────────────

def _parse_page_range(pages: str, total: int):
    if not pages:
        return 0, total
    parts = pages.split("-")
    start = int(parts[0]) - 1
    end   = int(parts[1]) if len(parts) > 1 else start + 1
    return max(0, start), min(end, total)


def _extract_pdf_pypdf(path: str, start: int, end: int):
    """用 pypdf 提取，回傳 (texts, meta) 或 (None, None) 若失敗"""
    try:
        import pypdf
        reader = pypdf.PdfReader(path)
        meta   = reader.metadata or {}
        texts  = []
        for i in range(start, end):
            txt = reader.pages[i].extract_text() or ""
            texts.append({"page": i + 1, "text": txt.strip()})
        return texts, meta
    except Exception:
        return None, None


def _extract_pdf_pdfplumber(path: str, start: int, end: int):
    """用 pdfplumber 提取（中文 PDF 備援），回傳 texts 或 None"""
    try:
        import pdfplumber
        texts = []
        with pdfplumber.open(path) as pdf:
            for i in range(start, end):
                if i >= len(pdf.pages):
                    break
                txt = pdf.pages[i].extract_text() or ""
                texts.append({"page": i + 1, "text": txt.strip()})
        return texts
    except Exception:
        return None


# Tesseract 可能的 Windows 安裝路徑
_TESSERACT_WIN_PATHS = [
    r"C:\Program Files\Tesseract-OCR\tesseract.exe",
    r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
]


def _find_tesseract() -> str | None:
    """回傳 tesseract 執行檔路徑，找不到回傳 None"""
    import shutil
    cmd = shutil.which("tesseract")
    if cmd:
        return cmd
    for p in _TESSERACT_WIN_PATHS:
        if os.path.exists(p):
            return p
    return None


def _ocr_pdf_pages(path: str, start: int, end: int) -> list | None:
    """
    掃描型 PDF OCR 提取。
    pypdfium2 渲染頁面（2× 縮放）→ Tesseract（chi_tra+eng）。
    回傳 texts 列表；Tesseract 未安裝時回傳 None。
    """
    tess_path = _find_tesseract()
    if not tess_path:
        return None

    # 渲染 PDF 頁面為 PIL Image
    try:
        import pypdfium2 as pdfium
        pdf_doc = pdfium.PdfDocument(path)
        page_images = []
        for i in range(start, min(end, len(pdf_doc))):
            page   = pdf_doc[i]
            bitmap = page.render(scale=2.0)   # 2× 提升 OCR 精度
            pil_img = bitmap.to_pil()
            page_images.append((i + 1, pil_img))
        pdf_doc.close()
    except Exception:
        return None

    if not page_images:
        return []

    # Tesseract OCR
    try:
        import pytesseract
        pytesseract.pytesseract.tesseract_cmd = tess_path

        # 偵測可用語言，優先用繁體中文
        available = pytesseract.get_languages(config="")
        if "chi_tra" in available:
            lang = "chi_tra+eng"
        elif "chi_sim" in available:
            lang = "chi_sim+eng"
        else:
            lang = "eng"

        texts = []
        for page_num, img in page_images:
            txt = pytesseract.image_to_string(
                img, lang=lang, config="--psm 3 --oem 3")
            texts.append({
                "page":       page_num,
                "text":       txt.strip(),
                "ocr_engine": f"tesseract/{lang}",
            })
        return texts
    except Exception:
        return None


def extract_pdf(path: str, pages: str = None) -> dict:
    try:
        import pypdf
        reader = pypdf.PdfReader(path)
        total  = len(reader.pages)
        meta   = reader.metadata or {}
    except ImportError:
        return {"error": "需要安裝 pypdf：pip install pypdf"}
    except Exception as e:
        return {"error": str(e)}

    start, end = _parse_page_range(pages, total)

    # 1. pypdf 提取
    texts, meta = _extract_pdf_pypdf(path, start, end)

    # 2. 若全空，改用 pdfplumber（中文 PDF 支援較佳）
    if texts is not None and all(not t["text"] for t in texts):
        plumber_texts = _extract_pdf_pdfplumber(path, start, end)
        if plumber_texts and any(t["text"] for t in plumber_texts):
            texts = plumber_texts

    if texts is None:
        texts = []

    all_empty = all(not t["text"] for t in texts)
    ocr_used  = False
    warning   = None

    # 3. 若仍全空，嘗試 OCR（掃描型 PDF）
    if all_empty:
        ocr_texts = _ocr_pdf_pages(path, start, end)
        if ocr_texts is not None:
            # OCR 有執行（不論結果是否為空）
            ocr_used = True
            if any(t["text"] for t in ocr_texts):
                texts = ocr_texts
                all_empty = False
            else:
                warning = "OCR 完成但未識別出文字，請確認文件影像品質。"
        else:
            # Tesseract 未安裝
            warning = (
                "此為掃描型 PDF，需要 OCR 才能提取文字。\n"
                "請安裝 Tesseract OCR（含繁體中文語言包 chi_tra）：\n"
                "https://github.com/UB-Mannheim/tesseract/wiki"
            )

    result = {
        "type":             "PDF",
        "total_pages":      total,
        "extracted_pages":  f"{start+1}-{min(end, total)}",
        "ocr_used":         ocr_used,
        "metadata": {
            "title":    meta.get("/Title",  "") if meta else "",
            "author":   meta.get("/Author", "") if meta else "",
            "creator":  meta.get("/Creator","") if meta else "",
            "created":  str(meta.get("/CreationDate", "")) if meta else "",
            "modified": str(meta.get("/ModDate",      "")) if meta else "",
        },
        "pages": texts,
    }
    if warning:
        result["warning"] = warning
    return result


# ── DOCX ──────────────────────────────────────────────────────────────────────

def extract_docx(path: str) -> dict:
    try:
        from docx import Document
        from docx.oxml.ns import qn
        doc  = Document(path)
        core = doc.core_properties

        # 段落
        paragraphs = []
        for p in doc.paragraphs:
            txt = p.text.strip()
            if txt:
                paragraphs.append({
                    "style": p.style.name if p.style else "",
                    "text":  txt,
                })

        # 表格
        tables = []
        for ti, tbl in enumerate(doc.tables):
            rows = []
            for row in tbl.rows:
                rows.append([c.text.strip() for c in row.cells])
            tables.append({"table_index": ti + 1, "rows": rows})

        return {
            "type": "DOCX",
            "metadata": {
                "title":    core.title or "",
                "author":   core.author or "",
                "created":  str(core.created or ""),
                "modified": str(core.modified or ""),
                "revision": core.revision or 0,
            },
            "paragraph_count": len(paragraphs),
            "table_count": len(tables),
            "paragraphs": paragraphs,
            "tables": tables,
        }
    except ImportError:
        return {"error": "需要安裝 python-docx：pip install python-docx"}
    except Exception as e:
        return {"error": str(e)}


# ── XLSX ──────────────────────────────────────────────────────────────────────

def extract_xlsx(path: str) -> dict:
    try:
        import openpyxl
        wb   = openpyxl.load_workbook(path, data_only=True)
        prop = wb.properties

        sheets = []
        for name in wb.sheetnames:
            ws   = wb[name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                row_data = [str(c) if c is not None else "" for c in row]
                if any(c for c in row_data):
                    rows.append(row_data)
            sheets.append({
                "name":      name,
                "row_count": ws.max_row,
                "col_count": ws.max_column,
                "data":      rows[:500],  # 最多 500 列
            })

        return {
            "type": "XLSX",
            "metadata": {
                "title":    prop.title or "",
                "creator":  prop.creator or "",
                "created":  str(prop.created or ""),
                "modified": str(prop.modified or ""),
            },
            "sheet_count": len(sheets),
            "sheets": sheets,
        }
    except ImportError:
        return {"error": "需要安裝 openpyxl：pip install openpyxl"}
    except Exception as e:
        return {"error": str(e)}


# ── PPTX ──────────────────────────────────────────────────────────────────────

def extract_pptx(path: str) -> dict:
    try:
        from pptx import Presentation
        prs  = Presentation(path)
        core = prs.core_properties

        slides = []
        for si, slide in enumerate(prs.slides):
            texts  = []
            images = 0
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        txt = para.text.strip()
                        if txt:
                            texts.append(txt)
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    images += 1
            slides.append({
                "slide": si + 1,
                "texts": texts,
                "image_count": images,
            })

        return {
            "type": "PPTX",
            "metadata": {
                "title":    core.title or "",
                "author":   core.author or "",
                "created":  str(core.created or ""),
                "modified": str(core.modified or ""),
            },
            "slide_count": len(slides),
            "slides": slides,
        }
    except ImportError:
        return {"error": "需要安裝 python-pptx：pip install python-pptx"}
    except Exception as e:
        return {"error": str(e)}


# ── ODT / ODS / ODP ──────────────────────────────────────────────────────────

def extract_odf(path: str) -> dict:
    try:
        from odf.opendocument import load
        from odf import text as odftext, table as odftable
        from odf.element import Element

        doc  = load(path)
        ext  = os.path.splitext(path)[1].lower()
        ftype = {"odt": "ODT", "ods": "ODS", "odp": "ODP"}.get(ext[1:], "ODF")

        # 提取所有文字節點
        def get_text(node) -> str:
            result = []
            if hasattr(node, "childNodes"):
                for child in node.childNodes:
                    if hasattr(child, "data"):
                        result.append(child.data)
                    else:
                        result.append(get_text(child))
            return "".join(result)

        body  = doc.body
        texts = []
        for elem in body.childNodes:
            t = get_text(elem).strip()
            if t:
                texts.append(t)

        meta = doc.meta
        def _meta(tag):
            try:
                return str(meta.getElementsByType(tag)[0]) if meta else ""
            except Exception:
                return ""

        return {
            "type": ftype,
            "metadata": {
                "title":    _meta("title"),
                "creator":  _meta("initial-creator"),
                "created":  _meta("creation-date"),
            },
            "text_blocks": len(texts),
            "content": texts,
        }
    except ImportError:
        return {"error": "需要安裝 odfpy：pip install odfpy"}
    except Exception as e:
        return {"error": str(e)}


# ── 圖片 ──────────────────────────────────────────────────────────────────────

def extract_image(path: str) -> dict:
    try:
        from PIL import Image
        from PIL.ExifTags import TAGS, GPSTAGS
        img  = Image.open(path)
        info: dict = {
            "type":   "IMAGE",
            "format": img.format or os.path.splitext(path)[1].upper()[1:],
            "mode":   img.mode,
            "width":  img.width,
            "height": img.height,
            "size_bytes": os.path.getsize(path),
        }

        # EXIF
        exif_data = {}
        raw_exif  = img._getexif() if hasattr(img, "_getexif") else None
        if raw_exif:
            for tag_id, value in raw_exif.items():
                tag = TAGS.get(tag_id, str(tag_id))
                if tag == "GPSInfo" and isinstance(value, dict):
                    gps = {}
                    for gk, gv in value.items():
                        gps[GPSTAGS.get(gk, gk)] = str(gv)
                    exif_data["GPSInfo"] = gps
                elif isinstance(value, bytes):
                    exif_data[tag] = value.hex()[:64]
                else:
                    try:
                        exif_data[tag] = str(value)[:256]
                    except Exception:
                        pass
        if exif_data:
            info["exif"] = exif_data

        # 檔案時間
        stat = os.stat(path)
        info["file_created"]  = datetime.datetime.fromtimestamp(
            stat.st_ctime).strftime("%Y-%m-%d %H:%M:%S")
        info["file_modified"] = datetime.datetime.fromtimestamp(
            stat.st_mtime).strftime("%Y-%m-%d %H:%M:%S")

        # GPS 座標換算
        if "GPSInfo" in exif_data:
            try:
                gps = exif_data["GPSInfo"]
                def dms_to_deg(dms_str):
                    import re
                    nums = re.findall(r"[\d.]+", dms_str)
                    return float(nums[0]) + float(nums[1])/60 + float(nums[2])/3600
                lat = dms_to_deg(str(gps.get("GPSLatitude", "0 0 0")))
                lon = dms_to_deg(str(gps.get("GPSLongitude", "0 0 0")))
                if gps.get("GPSLatitudeRef") == "S":  lat = -lat
                if gps.get("GPSLongitudeRef") == "W": lon = -lon
                info["gps_decimal"] = {"lat": round(lat, 7), "lon": round(lon, 7)}
            except Exception:
                pass

        return info
    except ImportError:
        return {"error": "需要安裝 Pillow：pip install Pillow"}
    except Exception as e:
        return {"error": str(e)}


# ── 主程式 ────────────────────────────────────────────────────────────────────

EXT_MAP = {
    ".pdf":  extract_pdf,
    ".docx": extract_docx,
    ".xlsx": extract_xlsx,
    ".pptx": extract_pptx,
    ".odt":  extract_odf,
    ".ods":  extract_odf,
    ".odp":  extract_odf,
    ".png":  extract_image,
    ".jpg":  extract_image,
    ".jpeg": extract_image,
    ".bmp":  extract_image,
    ".tiff": extract_image,
    ".tif":  extract_image,
    ".webp": extract_image,
    ".gif":  extract_image,
    ".heic": extract_image,
}

def main():
    parser = argparse.ArgumentParser(description="文件內容提取工具")
    parser.add_argument("file", help="檔案路徑")
    parser.add_argument("--pages", default=None, help="PDF 頁碼範圍（如 1-5）")
    parser.add_argument("--pretty", action="store_true", help="格式化 JSON 輸出")
    args = parser.parse_args()

    path = args.file
    if not os.path.exists(path):
        print(json.dumps({"error": f"檔案不存在：{path}"}))
        sys.exit(1)

    ext = os.path.splitext(path)[1].lower()
    fn  = EXT_MAP.get(ext)
    if not fn:
        print(json.dumps({
            "error": f"不支援的檔案類型：{ext}",
            "supported": list(EXT_MAP.keys())
        }))
        sys.exit(1)

    if ext == ".pdf" and args.pages:
        result = fn(path, pages=args.pages)
    else:
        result = fn(path)

    result["file_path"] = os.path.abspath(path)
    result["file_name"] = os.path.basename(path)
    result["file_size"] = os.path.getsize(path)

    indent = 2 if args.pretty else None
    output = json.dumps(result, ensure_ascii=False, indent=indent,
                        default=str)
    # 直接寫入 bytes，避免 stdout 編碼問題（CP950/Big5 環境）
    sys.stdout.buffer.write(output.encode("utf-8"))
    sys.stdout.buffer.write(b"\n")
    sys.stdout.buffer.flush()

if __name__ == "__main__":
    main()
